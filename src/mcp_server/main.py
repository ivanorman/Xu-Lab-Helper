# src/mcp_server/main.py
from __future__ import annotations
import os, json, hashlib, glob, subprocess, re, shutil
from typing import Iterator, Dict, List, Tuple, Set
from datetime import datetime

from mcp.server.fastmcp import FastMCP
from tqdm import tqdm

# ----------------------
# MCP instance
# ----------------------
mcp = FastMCP("xu-lab-decks")

# ----------------------
# Defaults / Paths
# ----------------------
DEF_ROOT   = os.environ.get("XU_DECKS_ROOT", "data/samples")
DEF_MAN    = os.environ.get("XU_DECKS_MANIFEST", "outputs/manifest.jsonl")
DEF_SLIDES = os.environ.get("XU_DECKS_SLIDES", "outputs/jsonl")
DEF_IMAGES = os.environ.get("XU_DECKS_IMAGES", "outputs/images")
DEF_SCHEMA = os.environ.get("XU_DECKS_SCHEMA", "outputs/schema_records")
DEF_SCHEMA_INDEX = os.environ.get("XU_DECKS_SCHEMA_INDEX", "outputs/index_schema")
SCHEMA_VERSION = "deck-schema-v1"
STATE_PATH = os.environ.get("XU_DECKS_STATE", "outputs/STATE.json")
ALIASES_PATH = os.environ.get("XU_DECKS_ALIASES", "outputs/aliases.json")

# ----------------------
# Helpers
# ----------------------

def _sha256_file(path: str, chunk: int = 1 << 20) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for b in iter(lambda: f.read(chunk), b""):
            h.update(b)
    return h.hexdigest()


def _walk_decks(root: str) -> Iterator[str]:
    for dirpath, _, files in os.walk(root):
        for name in files:
            if name.startswith((".", "~$")):
                continue
            low = name.lower()
            if low.endswith((".pptx", ".pdf")):
                yield os.path.join(dirpath, name)


def _listing_fingerprint(root: str) -> str:
    """Stable fingerprint of current source tree (.pptx/.pdf paths, sizes, mtimes)."""
    rows: List[str] = []
    for p in _walk_decks(root):
        st = os.stat(p)
        rows.append(f"{os.path.abspath(p)}|{st.st_size}|{int(st.st_mtime)}")
    payload = "\n".join(sorted(rows))
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def _read_state() -> Dict:
    if os.path.exists(STATE_PATH):
        try:
            with open(STATE_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {}


def _write_state(st: Dict) -> None:
    os.makedirs(os.path.dirname(STATE_PATH) or ".", exist_ok=True)
    st["last_run_iso"] = datetime.utcnow().isoformat(timespec="seconds") + "Z"
    with open(STATE_PATH, "w", encoding="utf-8") as f:
        json.dump(st, f, indent=2)

# ----------------------
# Date helpers
# ----------------------
from datetime import date

_DATE_PAT = re.compile(r"(20\d{2})(?:[._-]?(0[1-p]|1[0-2]))?(?:[._-]?([0-2]\d|3[01]))?")

def _infer_date_from_name(name: str) -> str | None:
    """Return ISO date string (YYYY-MM-DD or YYYY-MM or YYYY) if found in filename."""
    m = _DATE_PAT.search(name)
    if not m:
        return None
    y, mo, d = m.groups()
    if d:
        return f"{y}-{mo}-{d}"
    if mo:
        return f"{y}-{mo}"
    return y
def _iso_from_mtime(ts: float) -> str:
    return datetime.utcfromtimestamp(ts).date().isoformat()

def _parse_date_filter_from_query(q: str) -> tuple[str | None, str | None]:
    """Extract a coarse date range like '2023-2024' or a single year '2023'. Returns (from_iso, to_iso)."""
    # try ranges like 2023-2024, 2023/2024, 2023 to 2024
    rng = re.search(r"(20\d{2})\s*(?:[-/]|\\s+to\\s+)\s*(20\d{2})", q)
    if rng:
        y1, y2 = int(rng.group(1)), int(rng.group(2))
        if y1 > y2:
            y1, y2 = y2, y1
        return (f"{y1}-01-01", f"{y2}-12-31")
    yr = re.search(r"\b(20\d{2})\b", q)
    if yr:
        y = int(yr.group(1))
        return (f"{y}-01-01", f"{y}-12-31")
    return (None, None)

def _date_in_range(date_iso: str | None, start_iso: str | None, end_iso: str | None) -> bool:
    if not (start_iso or end_iso):
        return True
    if not date_iso:
        return False
    try:
        d = date.fromisoformat(date_iso if len(date_iso) == 10 else (date_iso + "-01-01" if len(date_iso) == 4 else date_iso + "-01"))
        ds = date.fromisoformat(start_iso) if start_iso else None
        de = date.fromisoformat(end_iso) if end_iso else None
        if ds and d < ds: return False
        if de and d > de: return False
        return True
    except Exception:
        return False


# ----------------------
# Alias registry helpers
# ----------------------
_WORD = re.compile(r"[\w\-µ°²³·]+", re.UNICODE)
_SUBSCRIPT_MAP = str.maketrans("₀₁₂₃₄₅₆₇₈₉", "0123456789")
def _asciiish(s: str) -> str:
    # normalize common unicode variants
    return (
        s.replace("–", "-").replace("—", "-")
         .replace("µ", "u").replace("·", " ")
         .translate(_SUBSCRIPT_MAP)
    )
def _canon_token(tok: str) -> str:
    t = _asciiish(tok).lower().strip()
    # compact formulas like "mo te2" -> "mote2"
    t = re.sub(r"\s+", " ", t)
    if re.fullmatch(r"[a-z]{1,3}\s*[a-z]{1,3}\s*\d{0,2}", t):
        t = t.replace(" ", "")
    return t
def _aliases_load() -> Dict[str, List[str]]:
    if os.path.exists(ALIASES_PATH):
        try:
            with open(ALIASES_PATH, "r", encoding="utf-8") as f:
                data = json.load(f)
                return { _canon_token(k): list(v) for k, v in data.items() }
        except Exception:
            pass
    return {}
def _aliases_save(d: Dict[str, List[str]]) -> None:
    os.makedirs(os.path.dirname(ALIASES_PATH) or ".", exist_ok=True)
    with open(ALIASES_PATH, "w", encoding="utf-8") as f:
        json.dump(d, f, indent=2, ensure_ascii=False)

def _twist_split(token: str) -> str | None:
    """If token looks like t<Material> (e.g., tMoTe2), return the remainder (e.g., MoTe2)."""
    if not token:
        return None
    if token.startswith("t") and len(token) > 1 and token[1].isupper():
        return token[1:]
    return None

def _expand_variants(token: str, reg: Dict[str, List[str]]) -> Set[str]:
    """Variants from registry + deterministic rules"""
    base = _canon_token(token)
    out: Set[str] = set()
    if base in reg:
        out.update(_asciiish(v).lower() for v in reg[base])
    out.add(_asciiish(token).lower())
    # tMaterial → treat as twisted shorthand
    if token.lower().startswith("t") and len(token) > 1 and token[1].isupper():
        out.add("t-")
    # chemical formulas spacing variants (also try after stripping a leading 't')
    form = _asciiish(token)
    form_core = form[1:] if (form.startswith("t") and len(form) > 1 and form[1].isupper()) else form
    m = re.fullmatch(r"([A-Za-z]{1,3})([A-Za-z]{1,3})(\d{0,2})", form_core)
    if m:
        a, b, n = m.groups()
        if n is None: n = ""
        out.add(f"{a}{b}{n}".lower())
        out.add(f"{a} {b}{n}".lower())
    return {v.strip() for v in out if v.strip()}

@mcp.tool("aliases_get", description="Return the current alias registry (canonical -> list of aliases)")
def aliases_get() -> Dict:
    return {"path": os.path.abspath(ALIASES_PATH), "aliases": _aliases_load()}

@mcp.tool("aliases_upsert", description="Upsert aliases for a canonical key; returns updated entry")
def aliases_upsert(canonical: str, aliases: List[str]) -> Dict:
    reg = _aliases_load()
    key = _canon_token(canonical)
    cur = set(reg.get(key, []))
    for a in aliases:
        if a and a.strip():
            cur.add(a.strip())
    reg[key] = sorted(cur)
    _aliases_save(reg)
    return {"canonical": key, "aliases": reg[key], "path": os.path.abspath(ALIASES_PATH)}



# ----------------------
# Refresh helper (idempotent): ensure corpus is ready
# ----------------------

def _prepare_corpus(root: str = DEF_ROOT, manifest: str = DEF_MAN, slides_dir: str = DEF_SLIDES, force: bool = False) -> Dict:
    """Check freshness; if stale, re-run ingest and parse. Returns final status."""
    st0 = status(root=root, manifest=manifest, slides_dir=slides_dir)
    if st0.get("is_fresh") and not force:
        st0["refreshed"] = False
        return st0
    # Rebuild deterministically
    ingest(root=root, out=manifest, force=force)
    parse(manifest=manifest, out=slides_dir, force=force)
    # Deterministically build schema and index after parsing
    draft_schema_all(slides_dir=slides_dir, out_dir=DEF_SCHEMA, force=force, template_md='schema.md')
    index_schema(records_dir=DEF_SCHEMA, out_dir=DEF_SCHEMA_INDEX, force=force)
    st1 = status(root=root, manifest=manifest, slides_dir=slides_dir)
    st1["refreshed"] = True
    return st1

@mcp.tool("ensure_ready", description="Ensure the slide corpus is up-to-date; rebuilds if needed and returns status")
def ensure_ready(root: str = DEF_ROOT, manifest: str = DEF_MAN, slides_dir: str = DEF_SLIDES, force: bool = False) -> Dict:
    return _prepare_corpus(root=root, manifest=manifest, slides_dir=slides_dir, force=force)

# ----------------------
# Tools
# ----------------------

@mcp.tool("status", description="Report whether manifest/slides are fresh; does not modify data")
def status(root: str = DEF_ROOT, manifest: str = DEF_MAN, slides_dir: str = DEF_SLIDES) -> Dict:
    """Return freshness and summary for current corpus outputs."""
    st = _read_state()

    # Compute current fingerprint of sources
    current_fp = _listing_fingerprint(root) if os.path.isdir(root) else ""

    manifest_exists = os.path.exists(manifest)
    slides_exist = os.path.isdir(slides_dir) and bool(glob.glob(os.path.join(slides_dir, "*.jsonl")))

    is_fresh = (
        manifest_exists
        and slides_exist
        and st.get("manifest_fp") == current_fp
        and st.get("slides_fp") == st.get("manifest_fp")
        and os.path.abspath(st.get("manifest_path", "")) == os.path.abspath(manifest)
        and os.path.abspath(st.get("slides_dir", "")) == os.path.abspath(slides_dir)
    )
    # Schema freshness
    schema_exist = os.path.isdir(DEF_SCHEMA) and bool(glob.glob(os.path.join(DEF_SCHEMA, "*.schema.json")))
    schema_fresh = schema_exist and st.get("schema_fp") == st.get("slides_fp") and st.get("schema_version") == SCHEMA_VERSION
    is_fresh = is_fresh and schema_fresh

    return {
        "source_root": os.path.abspath(root),
        "manifest": os.path.abspath(manifest),
        "slides_dir": os.path.abspath(slides_dir),
        "is_fresh": bool(is_fresh),
        "manifest_fp_recorded": st.get("manifest_fp"),
        "manifest_fp_current": current_fp,
        "slides_fp_recorded": st.get("slides_fp"),
        "slide_file_count": len(glob.glob(os.path.join(slides_dir, "*.jsonl"))) if os.path.isdir(slides_dir) else 0,
        "last_run": st.get("last_run_iso"),
        "schema_dir": os.path.abspath(DEF_SCHEMA),
        "schema_index": os.path.abspath(DEF_SCHEMA_INDEX),
        "schema_files": len(glob.glob(os.path.join(DEF_SCHEMA, "*.schema.json"))) if os.path.isdir(DEF_SCHEMA) else 0,
        "schema_fp_recorded": st.get("schema_fp"),
        "schema_version": st.get("schema_version"),
    }


@mcp.tool("ingest", description="Scan for .pptx/.pdf and write a JSONL manifest; idempotent with freshness check")
def ingest(root: str | None = None, out: str | None = None, force: bool = False) -> Dict:
    """
    Scan a folder for .pptx/.pdf and write a JSONL manifest.

    Defaults:
      root = $XU_DECKS_ROOT or "data/samples"
      out  = $XU_DECKS_MANIFEST or "outputs/manifest.jsonl"
    Set force=True to rebuild regardless of freshness.
    """
    root = root or DEF_ROOT
    out = out or DEF_MAN

    os.makedirs(os.path.dirname(out) or ".", exist_ok=True)

    # Freshness check
    current_fp = _listing_fingerprint(root)
    st = _read_state()
    if (st.get("manifest_fp") == current_fp) and os.path.exists(out) and not force:
        return {
            "skipped": True,
            "reason": "manifest up-to-date",
            "manifest": os.path.abspath(out),
            "files": sum(1 for _ in _walk_decks(root)),
        }

    files = list(_walk_decks(root))
    count = 0
    with open(out, "w", encoding="utf-8") as f:
        for p in tqdm(files, desc="Hashing"):
            st_os = os.stat(p)
            name = os.path.basename(p)
            inferred = _infer_date_from_name(name)
            mtime_iso = _iso_from_mtime(st_os.st_mtime)
            rec = {
                "path": os.path.abspath(p),
                "size": st_os.st_size,
                "mtime": st_os.st_mtime,
                "mtime_iso": mtime_iso,
                "sha256": _sha256_file(p),
                "date_iso": inferred,  # best-effort from filename
            }
            f.write(json.dumps(rec) + "\n")
            count += 1

    # record state
    st.update({
        "source_root": os.path.abspath(root),
        "manifest_path": os.path.abspath(out),
        "manifest_fp": current_fp,
    })
    _write_state(st)

    return {"skipped": False, "manifest": os.path.abspath(out), "files": count}


# ---------- PPTX parsing helper ----------

def _slide_text_from_pptx(path: str) -> List[str]:
    """Return a list[str] of per-slide extracted text for a .pptx file."""
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("python-pptx is not installed or failed to import") from e

    prs = Presentation(path)
    slide_texts: List[str] = []
    for slide in prs.slides:
        chunks: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                tf = shape.text_frame
                for para in tf.paragraphs:
                    line = "".join(run.text for run in para.runs) or para.text or ""
                    if line and line.strip():
                        chunks.append(line.strip())
        slide_texts.append("\n".join(chunks).strip())
    return slide_texts


@mcp.tool("parse", description="Parse PPTX decks from a manifest into per-slide JSONL; idempotent with freshness check")
def parse(manifest: str | None = None, out: str | None = None, force: bool = False, ocr: bool = False) -> Dict:
    """
    Parse PPTX decks from a manifest into per-slide JSONL files.

    Defaults:
      manifest = $XU_DECKS_MANIFEST or "outputs/manifest.jsonl"
      out      = $XU_DECKS_SLIDES or "outputs/jsonl"

    Input (manifest.jsonl): one line per file with {"path": "...", ...}
    Output: one file per deck → <out>/<basename>.jsonl
      Each line: {"file": {"path": ...}, "file_meta": {...}, "slide_index": i, "text": "...", "images": ["..."], "ocr": "..." (optional)}
    """
    manifest = manifest or DEF_MAN
    out_dir = out or DEF_SLIDES
    os.makedirs(out_dir, exist_ok=True)

    images_root = os.environ.get("XU_DECKS_IMAGES", DEF_IMAGES)
    os.makedirs(images_root, exist_ok=True)

    st = _read_state()
    # Skip only if slides were built for the same manifest fingerprint and exist
    if (
        st.get("slides_fp") == st.get("manifest_fp")
        and os.path.isdir(out_dir)
        and glob.glob(os.path.join(out_dir, "*.jsonl"))
        and not force
    ):
        return {
            "skipped": True,
            "reason": "slides up-to-date",
            "slides_dir": os.path.abspath(out_dir),
            "pptx_decks": len(glob.glob(os.path.join(out_dir, "*.jsonl"))),
        }

    decks = 0
    slides_total = 0
    manifest_meta: Dict[str, Dict] = {}
    seen_img_hashes: Set[str] = set()
    with open(manifest, "r", encoding="utf-8") as mf:
        for line in mf:
            if not line.strip():
                continue
            rec = json.loads(line)
            path = rec.get("path", "")
            manifest_meta[path] = {k: rec.get(k) for k in ("date_iso", "mtime_iso")}
            if not path.lower().endswith(".pptx"):
                # skip PDFs for now; PDF parser may be added later
                continue

            basename = os.path.splitext(os.path.basename(path))[0]
            out_path = os.path.join(out_dir, f"{basename}.jsonl")

            deck_img_dir = os.path.join(images_root, basename)
            os.makedirs(deck_img_dir, exist_ok=True)

            # Extract texts and embedded images per-slide
            try:
                from pptx import Presentation
                prs = Presentation(path)
                slide_texts: List[str] = []
                slide_images: List[List[str]] = []
                for s_idx, slide in enumerate(prs.slides):
                    # Text
                    chunks: List[str] = []
                    for shape in slide.shapes:
                        if getattr(shape, "has_text_frame", False):
                            tf = shape.text_frame
                            for para in tf.paragraphs:
                                line = "".join(run.text for run in para.runs) or para.text or ""
                                if line and line.strip():
                                    chunks.append(line.strip())
                    slide_texts.append("\n".join(chunks).strip())

                    # Images
                    cur_imgs: List[str] = []
                    for shape in slide.shapes:
                        # Picture shapes in python-pptx expose `.image`
                        img = getattr(shape, "image", None)
                        if img is None:
                            continue
                        try:
                            blob = img.blob
                            ihash = hashlib.sha256(blob).hexdigest()
                            if ihash in seen_img_hashes:
                                # de-dup identical images within this build
                                continue
                            seen_img_hashes.add(ihash)
                            ext = img.ext or "png"
                            fname = f"slide_{s_idx:03d}_img_{ihash[:8]}.{ext}"
                            ipath = os.path.join(deck_img_dir, fname)
                            # write atomically if possible
                            if not os.path.exists(ipath):
                                with open(ipath, "wb") as imf:
                                    imf.write(blob)
                            cur_imgs.append(os.path.abspath(ipath))
                        except Exception:
                            # ignore shapes that look like pictures but don't expose a blob
                            continue
                    slide_images.append(cur_imgs)
            except Exception as e:
                raise RuntimeError("python-pptx failed to parse images") from e

            # Optional OCR for each slide's images
            def _ocr_many(paths: List[str]) -> str:
                if not paths:
                    return ""
                try:
                    import pytesseract
                    from PIL import Image
                except Exception:
                    return ""  # silently skip if OCR stack is unavailable
                out_chunks: List[str] = []
                for pth in paths:
                    try:
                        txt = pytesseract.image_to_string(Image.open(pth))
                        t = (txt or "").strip()
                        if t:
                            out_chunks.append(t)
                    except Exception:
                        continue
                return "\n".join(out_chunks)

            with open(out_path, "w", encoding="utf-8") as outf:
                for i, txt in enumerate(slide_texts):
                    imgs = slide_images[i] if i < len(slide_images) else []
                    row = {
                        "file": {"path": os.path.abspath(path)},
                        "file_meta": manifest_meta.get(path, {}),
                        "slide_index": i,
                        "text": txt,
                        "images": imgs,
                    }
                    if ocr:
                        ocr_txt = _ocr_many(imgs)
                        if ocr_txt:
                            row["ocr"] = ocr_txt
                    outf.write(json.dumps(row, ensure_ascii=False) + "\n")

            decks += 1
            slides_total += len(slide_texts)

    # record state
    st.update({
        "slides_dir": os.path.abspath(out_dir),
        "slide_files": len(glob.glob(os.path.join(out_dir, "*.jsonl"))),
        # Record which manifest fingerprint these slides were built from
        "slides_fp": st.get("manifest_fp"),
    })
    _write_state(st)

    return {
        "skipped": False,
        "slides_dir": os.path.abspath(out_dir),
        "pptx_decks": decks,
        "slides": slides_total,
    }

# ----------------------
def _deck_image_dir(basename: str, images_root: str = DEF_IMAGES) -> str:
    return os.path.join(images_root, basename)

def _dir_size(path: str) -> int:
    total = 0
    if not os.path.isdir(path):
        return 0
    for dirpath, _, filenames in os.walk(path):
        for fn in filenames:
            fp = os.path.join(dirpath, fn)
            try:
                total += os.path.getsize(fp)
            except Exception:
                continue
    return total

# ----------------------
# Image extraction and OCR tools
# ----------------------


@mcp.tool("ocr_images", description="Run OCR on a list of image paths using Tesseract; returns {path: text}")
def ocr_images(paths: List[str]) -> Dict:
    try:
        import pytesseract
        from PIL import Image
    except Exception as e:
        return {"error": f"OCR stack unavailable: {e}"}
    out = {}
    for p in paths or []:
        try:
            txt = pytesseract.image_to_string(Image.open(p))
            if txt and txt.strip():
                out[os.path.abspath(p)] = txt.strip()
        except Exception:
            continue
    return {"results": out, "count": len(out)}


# List exported images for a deck and report total size in bytes
@mcp.tool("list_deck_images", description="List exported images for a deck and report total size in bytes")
def list_deck_images(basename: str, images_root: str = DEF_IMAGES) -> Dict:
    img_dir = _deck_image_dir(basename, images_root)
    files: List[str] = []
    if os.path.isdir(img_dir):
        for dirpath, _, fns in os.walk(img_dir):
            for fn in fns:
                files.append(os.path.abspath(os.path.join(dirpath, fn)))
    return {
        "basename": basename,
        "image_dir": os.path.abspath(img_dir),
        "count": len(files),
        "size_bytes": _dir_size(img_dir),
        "files": files,
    }

# Prune a deck's exported images: delete or zip-archive then delete
@mcp.tool("prune_deck_images", description="Prune a deck's exported images: delete or zip-archive then delete")
def prune_deck_images(basename: str, mode: str = "delete", images_root: str = DEF_IMAGES, zip_out_dir: str = "outputs/pruned") -> Dict:
    img_dir = _deck_image_dir(basename, images_root)
    if not os.path.isdir(img_dir):
        return {
            "basename": basename,
            "pruned": 0,
            "reason": "no image directory",
            "image_dir": os.path.abspath(img_dir),
        }
    before = list_deck_images(basename, images_root)
    archived_path = None
    if mode == "zip":
        os.makedirs(zip_out_dir, exist_ok=True)
        base = os.path.join(zip_out_dir, f"{basename}_images")
        try:
            archive = shutil.make_archive(base, "zip", root_dir=img_dir)
            archived_path = os.path.abspath(archive)
        except Exception as e:
            return {"error": f"zip failed: {e}", "image_dir": os.path.abspath(img_dir)}
    try:
        shutil.rmtree(img_dir)
    except Exception as e:
        return {"error": f"delete failed: {e}", "image_dir": os.path.abspath(img_dir)}
    return {
        "basename": basename,
        "image_dir_deleted": os.path.abspath(img_dir),
        "archived_zip": archived_path,
        "removed_count": before.get("count", 0),
        "freed_bytes": before.get("size_bytes", 0),
        "mode": mode,
    }


# ----------------------
# Search with synonym expansion
# ----------------------
def _token_groups_from_query(q: str, reg: Dict[str, List[str]]) -> List[Set[str]]:
    raw = [t for t in _WORD.findall(q) if t]
    groups: List[Set[str]] = []
    for tok in raw:
        remainder = _twist_split(tok)
        if remainder:
            # Require both concepts: twisting AND the material.
            # Also allow the compact form (e.g., "tMoTe2") to satisfy BOTH groups
            # by including the original token in each set.
            tw = _expand_variants("twisted", reg)
            mat = _expand_variants(remainder, reg)
            orig = tok.lower()
            tw.add(orig)
            mat.add(orig)
            groups.append(tw)
            groups.append(mat)
        else:
            groups.append(_expand_variants(tok, reg))
    groups = [g for g in groups if g]
    return groups

def _match_groups(text: str, groups: List[Set[str]]) -> Tuple[bool, int]:
    t = _asciiish(text).lower()
    score = 0
    for g in groups:
        group_hits = sum(t.count(v) for v in g if v)
        if group_hits == 0:
            return (False, 0)
        score += group_hits
    return (True, score)

@mcp.tool("search_slides", description="Search parsed slide JSONLs with synonym/alias expansion; returns ranked matches")
def search_slides(q: str, slides_dir: str = DEF_SLIDES, k: int = 25) -> Dict:
    reg = _aliases_load()
    # Guardrail: surface staleness to the client so it can decide to rebuild
    st = status(slides_dir=slides_dir)
    if not st.get("is_fresh"):
        return {"results": [], "total": 0, "query": q, "slides_dir": os.path.abspath(slides_dir), "needs_refresh": True, "status": st}
    start_iso, end_iso = _parse_date_filter_from_query(q)
    groups = _token_groups_from_query(q, reg)
    if not groups:
        return {"results": [], "error": "empty or un-tokenizable query", "query": q}
    if not os.path.isdir(slides_dir):
        return {"results": [], "error": f"slides_dir not found: {slides_dir}", "query": q}

    hits: List[Dict] = []
    for jf in glob.glob(os.path.join(slides_dir, "*.jsonl")):
        with open(jf, "r", encoding="utf-8") as f:
            for line in f:
                if not line.strip():
                    continue
                row = json.loads(line)
                meta = row.get("file_meta", {})
                file_date = meta.get("date_iso") or meta.get("mtime_iso")
                if not _date_in_range(file_date, start_iso, end_iso):
                    continue
                txt = row.get("text", "")
                matched, score = _match_groups(txt, groups)
                if matched:
                    anchor = next(iter(groups[0]))
                    tnorm = _asciiish(txt).lower()
                    i = tnorm.find(anchor)
                    if i == -1: i = 0
                    start = max(0, i - 120); end = min(len(txt), i + 120)
                    snippet = ("…" if start > 0 else "") + txt[start:end] + ("…" if end < len(txt) else "")
                    hits.append({
                        "file": row.get("file", {}).get("path"),
                        "slide_index": row.get("slide_index"),
                        "score": int(score),
                        "snippet": snippet
                    })
    hits.sort(key=lambda x: (-x["score"], x.get("file") or "", int(x.get("slide_index") or 0)))
    k = max(1, int(k))
    return {
        "results": hits[:k],
        "total": len(hits),
        "query": q,
        "slides_dir": os.path.abspath(slides_dir),
        "date_from": start_iso,
        "date_to": end_iso,
    }


@mcp.tool("open_presentation", description="Open a presentation file in the system viewer")
def open_presentation(path: str):
    """Opens a PowerPoint or PDF file using the system's default application."""
    subprocess.Popen(["open", path])
    return {"status": "opened", "path": path}

# ----------------------
# Slide row helpers and schema extraction
# ----------------------
def _iter_deck_rows(slides_dir: str, basename: str) -> Iterator[Dict]:
    jf = os.path.join(slides_dir, f"{basename}.jsonl")
    if not os.path.exists(jf):
        return
    with open(jf, "r", encoding="utf-8") as f:
        for line in f:
            if not line.strip():
                continue
            yield json.loads(line)

_MAT_PAT = re.compile(r"\b([A-Z][a-z]?[A-Z][a-z]?\d{0,2})\b")
_MEAS_MAP = {
    "rmcd": "rmcd",
    "mcd": "rmcd",
    "mr": "magnetoresistance",
    "magnetoresistance": "magnetoresistance",
    "transport": "transport",
    "iv": "iv",
    "rxx": "transport",
    "reflectance": "reflectance",
    "pl": "photoluminescence",
    "photoluminescence": "photoluminescence",
    "raman": "raman",
}
_NUM_PAT = re.compile(r"(?P<val>[+-]?[0-9]*\.?[0-9]+)\s*(?P<Unit>K|T|V\/nm|V\\/nm|cm^-2|cm^\-2|cm-2)")

def _derive_schema_from_rows(rows: List[Dict]) -> Dict:
    texts = []
    mats: Set[str] = set()
    meas: Set[str] = set()
    t_min = None; t_max = None
    b_min = None; b_max = None
    any_date = None
    for r in rows:
        meta = r.get("file_meta", {})
        any_date = any_date or meta.get("date_iso") or meta.get("mtime_iso")
        t = (r.get("text") or "") + "\n" + (r.get("ocr") or "")
        # Incorporate LLM image captions if present
        captions = r.get("image_captions")
        if captions and isinstance(captions, list):
            caption_texts = [c.get("caption", "") for c in captions if isinstance(c, dict) and c.get("caption")]
            if caption_texts:
                t = t + "\n" + "\n".join(caption_texts)
        t_norm = _asciiish(t).lower()
        texts.append(t)
        # materials: heuristic chemical tokens
        for m in _MAT_PAT.findall(t):
            mats.add(m)
        # measurement types
        for k, v in _MEAS_MAP.items():
            if k in t_norm:
                meas.add(v)
        # numeric cues
        for m in _NUM_PAT.finditer(t):
            val = None
            try:
                val = float(m.group("val"))
            except Exception:
                continue
            unit = m.group("Unit")
            if unit == "K":
                t_min = val if t_min is None else min(t_min, val)
                t_max = val if t_max is None else max(t_max, val)
            elif unit == "T":
                b_min = val if b_min is None else min(b_min, val)
                b_max = val if b_max is None else max(b_max, val)
    summary_hint = " ".join(sorted(list(meas))[:3])
    schema = {
        "schema_version": SCHEMA_VERSION,
        "date": any_date,
        "materials": sorted(mats) if mats else None,
        "experiment_type": sorted(list(meas)) if meas else None,
        "temperature_K_range": [t_min, t_max] if (t_min is not None and t_max is not None) else None,
        "magnetic_field_T_range": [b_min, b_max] if (b_min is not None and b_max is not None) else None,
        "summary": summary_hint or None,
        "notes": None,
    }
    # drop Nones for compactness
    return {k: v for k, v in schema.items() if v is not None}

# ----------------------
# Markdown schema template helpers
# ----------------------

_AUTO_FIELDS: Set[str] = {
    "schema_version",
    "date",
    "materials",
    "experiment_type",
    "temperature_K_range",
    "magnetic_field_T_range",
    "summary",
    "file",
    "provenance",
}

_MD_JSON_FENCE = re.compile(r"```json(.*?)```", re.DOTALL | re.IGNORECASE)
_MD_FIELD_LINE = re.compile(r"^\s*[-*]\s*(?:\*\*|`)?\s*([A-Za-z0-9_./ -]+?)\s*(?:\*\*|`)?:")
_MD_HEADING = re.compile(r"^\s*#{2,6}\s+([A-Za-z0-9_./ -]+?)\s*$")


def _read_text(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def _field_key(name: str) -> str:
    # convert human heading like "Key Findings" -> key_findings
    k = name.strip().lower()
    k = re.sub(r"[^a-z0-9]+", "_", k).strip("_")
    return k or name.strip()


def _parse_template_md(md_path: str) -> Dict:
    """Parse a Markdown template to extract an ordered list of field names.
    Strategies (in order):
      1) A fenced ```json block – parse keys from JSON object.
      2) Bullet lines like "- Field Name:" -> field_name
      3) Headings (##/### etc.) -> field_name
    Returns {"fields_order": [...], "fields": {key: None, ...}}
    """
    try:
        txt = _read_text(md_path)
    except Exception as e:
        return {"error": f"failed to read template: {e}", "fields_order": [], "fields": {}}

    # 1) JSON fence
    m = _MD_JSON_FENCE.search(txt)
    if m:
        block = m.group(1)
        try:
            obj = json.loads(block)
            if isinstance(obj, dict):
                order = list(obj.keys())
                return {"fields_order": order, "fields": {k: None for k in order}}
        except Exception:
            pass

    # 2) Bullet lines
    fields: List[str] = []
    for line in txt.splitlines():
        mb = _MD_FIELD_LINE.match(line)
        if mb:
            fields.append(_field_key(mb.group(1)))
    if fields:
        return {"fields_order": fields, "fields": {k: None for k in fields}}

    # 3) Headings
    for line in txt.splitlines():
        mh = _MD_HEADING.match(line)
        if mh:
            fields.append(_field_key(mh.group(1)))
    fields = [f for f in fields if f]
    return {"fields_order": fields, "fields": {k: None for k in fields}}


@mcp.tool("load_schema_template", description="Read a Markdown schema template and return detected field names in order")
def load_schema_template(md_path: str) -> Dict:
    info = _parse_template_md(md_path)
    info["path"] = os.path.abspath(md_path)
    return info


# Utility to assemble rich text context for NL fields

def _assemble_deck_context(rows: List[Dict], max_chars: int = 8000) -> str:
    parts: List[str] = []
    for r in rows:
        sidx = r.get("slide_index")
        header = f"[Slide {sidx}]"
        txt = r.get("text") or ""
        ocr = r.get("ocr") or ""
        caps = "\n".join(
            [c.get("caption", "") for c in (r.get("image_captions") or []) if isinstance(c, dict)]
        )
        blob = "\n".join([p for p in [header, txt, ocr, caps] if p]).strip()
        if blob:
            parts.append(blob)
        if sum(len(p) for p in parts) > max_chars:
            break
    ctx = "\n\n".join(parts)
    if len(ctx) > max_chars:
        ctx = ctx[:max_chars] + "…"
    return ctx

# ----------------------
# Schema tools
# ----------------------

@mcp.tool("build_schema_record", description="Build a deck-level schema JSON for a single deck basename (without extension)")
def build_schema_record(basename: str, slides_dir: str = DEF_SLIDES, out_dir: str = DEF_SCHEMA) -> Dict:
    os.makedirs(out_dir, exist_ok=True)
    rows = list(_iter_deck_rows(slides_dir, basename) or [])
    if not rows:
        return {"error": f"no slides for {basename}"}
    rec = _derive_schema_from_rows(rows)
    # extract file path from first row for provenance
    file_path = rows[0].get("file", {}).get("path")
    rec["file"] = {"path": file_path}
    # include simple provenance: slide count
    rec["provenance"] = {"slides": len(rows)}
    out_path = os.path.join(out_dir, f"{basename}.schema.json")
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(rec, f, indent=2, ensure_ascii=False)
    return {"basename": basename, "schema_path": os.path.abspath(out_path)}

@mcp.tool("build_schema_all", description="Build schema JSON for all decks in slides_dir; idempotent vs slides fingerprint")
def build_schema_all(slides_dir: str = DEF_SLIDES, out_dir: str = DEF_SCHEMA, force: bool = False) -> Dict:
    os.makedirs(out_dir, exist_ok=True)
    st = _read_state()
    slides_fp = st.get("slides_fp")
    # Skip if schema already built for current slides and not forced
    if st.get("schema_fp") == slides_fp and st.get("schema_version") == SCHEMA_VERSION and not force:
        return {"skipped": True, "reason": "schema up-to-date", "schema_dir": os.path.abspath(out_dir)}
    built = 0
    for jf in glob.glob(os.path.join(slides_dir, "*.jsonl")):
        basename = os.path.splitext(os.path.basename(jf))[0]
        res = build_schema_record(basename=basename, slides_dir=slides_dir, out_dir=out_dir)
        if not res.get("error"):
            built += 1
    # record state
    st.update({
        "schema_dir": os.path.abspath(out_dir),
        "schema_fp": slides_fp,
        "schema_version": SCHEMA_VERSION,
        "schema_files": len(glob.glob(os.path.join(out_dir, "*.schema.json"))),
    })
    _write_state(st)
    return {"skipped": False, "built": built, "schema_dir": os.path.abspath(out_dir)}

# ----------------------
# Template-driven schema tools
# ----------------------

@mcp.tool("draft_schema_record", description="Create/overwrite a deck schema JSON shaped by a Markdown template; auto-fill what code can, leave NL fields empty, and include a summarization context for the client.")
def draft_schema_record(basename: str, template_md: str, slides_dir: str = DEF_SLIDES, out_dir: str = DEF_SCHEMA, force: bool = False) -> Dict:
    os.makedirs(out_dir, exist_ok=True)
    rows = list(_iter_deck_rows(slides_dir, basename) or [])
    if not rows:
        return {"error": f"no slides for {basename}"}

    # Load template fields to define the shape of the blank schema
    ti = _parse_template_md(template_md)
    order = ti.get("fields_order") or []
    if not order:
        return {"error": "template yielded no fields", "template": os.path.abspath(template_md)}

    # Create a blank record with null values for all fields from the template.
    rec: Dict[str, object] = {k: None for k in order}

    # Add only essential, non-heuristic metadata.
    rec["schema_version"] = SCHEMA_VERSION
    rec["file"] = {"path": rows[0].get("file", {}).get("path")}
    rec["provenance"] = {"slides": len(rows)}

    # Attach the summarization context for the client LLM, which is the primary input for the next step.
    ctx = _assemble_deck_context(rows)
    rec["_nl_context"] = ctx

    out_path = os.path.join(out_dir, f"{basename}.schema.json")
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(rec, f, indent=2, ensure_ascii=False)
    return {"basename": basename, "schema_path": os.path.abspath(out_path), "template": os.path.abspath(template_md)}


@mcp.tool("draft_schema_all", description="Generate/refresh schema drafts for all decks using a Markdown template")
def draft_schema_all(template_md: str, slides_dir: str = DEF_SLIDES, out_dir: str = DEF_SCHEMA, force: bool = False) -> Dict:
    built = 0
    errs: List[str] = []
    for jf in glob.glob(os.path.join(slides_dir, "*.jsonl")):
        basename = os.path.splitext(os.path.basename(jf))[0]
        res = draft_schema_record(basename=basename, template_md=template_md, slides_dir=slides_dir, out_dir=out_dir, force=force)
        if res.get("error"):
            errs.append(f"{basename}: {res['error']}")
        else:
            built += 1
    return {"built": built, "errors": errs, "schema_dir": os.path.abspath(out_dir)}


@mcp.tool("commit_schema_fields", description="Update a deck schema JSON with user/LLM-provided natural-language fields; preserves auto fields")
def commit_schema_fields(basename: str, updates: Dict, out_dir: str = DEF_SCHEMA) -> Dict:
    fp = os.path.join(out_dir, f"{basename}.schema.json")
    if not os.path.exists(fp):
        return {"error": f"schema not found for {basename}", "schema_path": os.path.abspath(fp)}
    try:
        with open(fp, "r", encoding="utf-8") as f:
            rec = json.load(f)
    except Exception as e:
        return {"error": f"failed to read schema: {e}", "schema_path": os.path.abspath(fp)}

    if not isinstance(updates, dict):
        return {"error": "updates must be a JSON object"}

    for k, v in updates.items():
        if k not in _AUTO_FIELDS:
            rec[k] = v

    try:
        with open(fp, "w", encoding="utf-8") as f:
            json.dump(rec, f, indent=2, ensure_ascii=False)
    except Exception as e:
        return {"error": f"failed to write updated schema: {e}", "schema_path": os.path.abspath(fp)}

    return {"basename": basename, "schema_path": os.path.abspath(fp), "updated_fields": list(updates.keys())}


# ----------------------
# Indexing and searching schema
# ----------------------

def _index_batch(records: List[Dict], writer):
    from whoosh.fields import TEXT, ID, KEYWORD, NUMERIC
    # Add documents to the writer
    for r in records:
        path = r.get("file", {}).get("path")
        if not path: continue
        # Ranges need to be stored as separate fields for Whoosh
        temp_range = r.get("temperature_K_range")
        t_min, t_max = (None, None)
        if isinstance(temp_range, list) and len(temp_range) == 2:
            t_min, t_max = temp_range
        field_range = r.get("magnetic_field_T_range")
        b_min, b_max = (None, None)
        if isinstance(field_range, list) and len(field_range) == 2:
            b_min, b_max = field_range
        writer.add_document(
            path=path,
            basename=os.path.splitext(os.path.basename(path))[0],
            materials=r.get("materials"),
            experiment_type=r.get("experiment_type"),
            summary=r.get("summary"),
            notes=r.get("notes"),
            date=r.get("date"),
            temp_min_k=t_min,
            temp_max_k=t_max,
            field_min_t=b_min,
            field_max_t=b_max,
        )

@mcp.tool("index_schema", description="Index deck-level schema records for fast search (batched and safe).")
def index_schema(records_dir: str = DEF_SCHEMA, out_dir: str = DEF_SCHEMA_INDEX, force: bool = False, batch_size: int = 128) -> Dict:
    from whoosh.index import create_in, open_dir
    from whoosh.fields import Schema, TEXT, ID, KEYWORD, NUMERIC
    from whoosh.writing import AsyncWriter

    os.makedirs(out_dir, exist_ok=True)
    schema = Schema(
        path=ID(stored=True, unique=True),
        basename=ID(stored=True),
        materials=KEYWORD(stored=True, commas=True, scorable=True),
        experiment_type=KEYWORD(stored=True, commas=True, scorable=True),
        summary=TEXT(stored=True),
        notes=TEXT(stored=True),
        date=ID(stored=True),
        temp_min_k=NUMERIC(float, stored=True),
        temp_max_k=NUMERIC(float, stored=True),
        field_min_t=NUMERIC(float, stored=True),
        field_max_t=NUMERIC(float, stored=True),
    )
    st = _read_state()
    slides_fp = st.get("slides_fp")
    if not force and st.get("schema_index_fp") == slides_fp and st.get("schema_version") == SCHEMA_VERSION:
        return {"skipped": True, "reason": "index is up-to-date", "index_dir": os.path.abspath(out_dir)}

    ix = create_in(out_dir, schema)
    writer = AsyncWriter(ix)
    records = []
    for jf in glob.glob(os.path.join(records_dir, "*.schema.json")):
        with open(jf, "r", encoding="utf-8") as f:
            records.append(json.load(f))
        if len(records) >= batch_size:
            _index_batch(records, writer)
            records = []
    if records:
        _index_batch(records, writer)
    writer.commit()
    return {"indexed": ix.doc_count(), "index_dir": os.path.abspath(out_dir)}


@mcp.tool("index_schema_one", description="Update the schema index for a single deck by basename; no full rebuild.")
def index_schema_one(basename: str, records_dir: str = DEF_SCHEMA, out_dir: str = DEF_SCHEMA_INDEX) -> Dict:
    from whoosh.index import open_dir
    from whoosh.writing import AsyncWriter
    if not os.path.exists(out_dir):
        return {"error": "index does not exist; run index_schema first"}
    ix = open_dir(out_dir)
    writer = AsyncWriter(ix)
    jf = os.path.join(records_dir, f"{basename}.schema.json")
    if not os.path.exists(jf):
        return {"error": f"schema record not found for {basename}"}
    with open(jf, "r", encoding="utf-8") as f:
        record = json.load(f)
    _index_batch([record], writer)
    writer.commit()
    return {"basename": basename, "indexed": True, "index_dir": os.path.abspath(out_dir)}


@mcp.tool("search_schema", description="Search deck-level schema index with alias expansion and optional date parsing")
def search_schema(q: str, index_dir: str = DEF_SCHEMA_INDEX, k: int = 10) -> Dict:
    from whoosh.index import open_dir
    from whoosh.qparser import QueryParser, MultifieldParser
    from whoosh.query import Term, And, Or

    if not os.path.isdir(index_dir) or not os.path.exists(os.path.join(index_dir, "_MAIN_WRITELOCK")):
        return {"error": f"index not found at {index_dir}", "results": []}

    ix = open_dir(index_dir)
    reg = _aliases_load()
    groups = _token_groups_from_query(q, reg)
    
    # Each group is a set of synonyms; documents must match at least one term from each group.
    # e.g., "twisted MoTe2" -> (twisted OR t-mote2) AND (mote2 OR mo te2)
    query_parts = []
    for g in groups:
        if len(g) == 1:
            term = next(iter(g))
            query_parts.append(Or([Term("summary", term), Term("notes", term), Term("materials", term), Term("experiment_type", term)]))
        else:
            query_parts.append(Or([Term("summary", t) for t in g] + [Term("notes", t) for t in g] + [Term("materials", t) for t in g] + [Term("experiment_type", t) for t in g]))

    final_query = And(query_parts) if query_parts else None

    results = []
    if final_query:
        with ix.searcher() as searcher:
            res = searcher.search(final_query, limit=k)
            for hit in res:
                results.append({
                    "path": hit.get("path"),
                    "basename": hit.get("basename"),
                    "score": hit.score,
                    "date": hit.get("date"),
                    "summary": hit.get("summary"),
                    "materials": hit.get("materials"),
                    "experiment_type": hit.get("experiment_type"),
                })

    return {"query": q, "expanded_query": str(final_query), "results": results, "count": len(results)}

# ----------------------
# LLM-assisted tools
# ----------------------

@mcp.tool("prepare_caption_jobs", description="Prepare image-caption jobs for a deck; returns prompts and image paths for the client LLM to caption")
def prepare_caption_jobs(basename: str, slides_dir: str = DEF_SLIDES, max_images_per_slide: int = 5) -> Dict:
    rows = list(_iter_deck_rows(slides_dir, basename) or [])
    if not rows:
        return {"error": f"no slides for {basename}"}
    jobs = []
    for r in rows:
        sidx = r.get("slide_index")
        text = r.get("text") or ""
        imgs = r.get("images") or []
        if not imgs:
            continue
        # For now, a generic prompt. Could be improved.
        prompt = f"The following slide (index {sidx}) has this text:\\n---\\n{text}\\n---\\nDescribe the key features and data presented in the following image(s) from this slide. Focus on quantitative details if possible."
        # Limit images to avoid overwhelming the context window
        for i, img_path in enumerate(imgs[:max_images_per_slide]):
            jobs.append({
                "slide_index": sidx,
                "figure_id": f"slide_{sidx}_img_{i}",
                "prompt": prompt,
                "image_path": img_path,
            })
    return {"basename": basename, "jobs": jobs, "count": len(jobs)}


@mcp.tool("commit_captions", description="Write LLM-generated captions back into the deck JSONL. Input is a list of {slide_index, figure_id, caption, entities?, numbers?}.")
def commit_captions(basename: str, captions: List[Dict], slides_dir: str = DEF_SLIDES) -> Dict:
    jf = os.path.join(slides_dir, f"{basename}.jsonl")
    if not os.path.exists(jf):
        return {"error": f"no slides for {basename}"}
    
    # Group captions by slide index for efficient updates
    captions_by_slide: Dict[int, List[Dict]] = {}
    for cap in captions:
        sidx = cap.get("slide_index")
        if sidx is not None:
            if sidx not in captions_by_slide:
                captions_by_slide[sidx] = []
            captions_by_slide[sidx].append(cap)

    rows = list(_iter_deck_rows(slides_dir, basename) or [])
    updated_rows = 0
    with open(jf, "w", encoding="utf-8") as f:
        for r in rows:
            sidx = r.get("slide_index")
            if sidx in captions_by_slide:
                if "image_captions" not in r:
                    r["image_captions"] = []
                r["image_captions"].extend(captions_by_slide[sidx])
                updated_rows +=1
            f.write(json.dumps(r) + "\n")

    return {"basename": basename, "jsonl_path": os.path.abspath(jf), "updated_slides": updated_rows}

# ----------------------
# Smart search (schema preferred, fallback to slides)
# ----------------------

@mcp.tool("smart_search", description="Ensure corpus freshness, then search schema (preferred) with fallback to slide search")
def smart_search(q: str, root: str = DEF_ROOT, manifest: str = DEF_MAN, slides_dir: str = DEF_SLIDES, k: int = 10) -> Dict:
    # Always ensure corpus is fresh before searching
    fresh_status = _prepare_corpus(root=root, manifest=manifest, slides_dir=slides_dir, force=False)
    
    # Try schema search first
    schema_results = search_schema(q=q, k=k)
    if not schema_results.get("error") and schema_results.get("results"):
        return {
            "search_type": "schema",
            "results": schema_results["results"],
            "query": q,
            "fresh_status": fresh_status,
        }
    
    # Fallback to slide text search
    slide_results = search_slides(q=q, slides_dir=slides_dir, k=k)
    # Don't surface needs_refresh error from fallback, as _prepare_corpus already ran
    if "needs_refresh" in slide_results:
        del slide_results["needs_refresh"]

    return {
        "search_type": "slides",
        "results": slide_results.get("results", []),
        "query": q,
        "fresh_status": fresh_status,
    }


# ----------------------
# Context extraction for NL summarization
# ----------------------

@mcp.tool("extract_context", description="Return concatenated slide text, OCR and existing image captions for a deck for use in NL summarization")
def extract_context(basename: str, slides_dir: str = DEF_SLIDES, max_chars: int = 12000) -> Dict:
    rows = list(_iter_deck_rows(slides_dir, basename) or [])
    if not rows:
        return {"error": f"no slides for {basename}"}
    
    ctx = _assemble_deck_context(rows, max_chars)
    return {
        "basename": basename,
        "context": ctx,
        "chars": len(ctx),
    }

if __name__ == "__main__":
    mcp.run()
